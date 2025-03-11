import os
from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.responses import JSONResponse
from fastapi.staticfiles import StaticFiles
from typing import Dict, Any, Optional
import json
import logging
from datetime import datetime
import asyncio
import re

from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import tempfile

from fastapi.staticfiles import StaticFiles


# utils에서 필요한 클래스들 import
from utils.dc import DocumentProcessor
from utils.pg import ProposalGenerator

# 로깅 설정
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()

# 정적 파일 설정
app.mount("/static", StaticFiles(directory="static"), name="static")


# HTML 파일용 라우트 추가
@app.get("/static/html/prop.html")
async def get_prop_html():
    return FileResponse("static/html/prop.html")

# 제안서 상태 관리를 위한 Enum 클래스 추가
from enum import Enum
class ProposalStatus(str, Enum):
    INITIALIZED = "initialized"
    FILE_UPLOADED = "file_uploaded"
    ANALYSIS_COMPLETED = "analysis_completed"
    PLANNING_IN_PROGRESS = "planning_in_progress"
    PLANNING_COMPLETED = "planning_completed"
    GENERATING = "generating"
    COMPLETED = "completed"
    ERROR = "error"

class ProposalManager:
    def __init__(self):
        self.storage_dir = "proposals"
        self.ensure_storage_directory()
        self.doc_processor = DocumentProcessor()
        self.proposal_generator = ProposalGenerator()
        self.active_proposals: Dict[str, Dict] = {}

    def ensure_storage_directory(self):
        if not os.path.exists(self.storage_dir):
            os.makedirs(self.storage_dir)

    async def create_proposal(self, file_content: bytes, filename: str) -> str:
        try:
            # 새 제안서 ID 생성
            proposal_id = f"proposal_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            
            # 파일 처리 및 초기 분석
            doc, page_limit = self.doc_processor.process_file(file_content)
            
            if doc is None:
                raise HTTPException(status_code=400, detail="파일 처리 실패")
                
            # 파일이 hwp인 경우
            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext in ['.hwp', '.hwpx']:
                result = self.doc_processor._process_hwp(file_content)
                if result is None:
                    raise HTTPException(status_code=400, detail="HWP 파일 처리 실패")
                full_text = result.get('full_text', '')
                table_sections = result.get('table_sections', [])
            else:  # PDF 파일의 경우
                full_text = ""
                for page in doc.pages if hasattr(doc, 'pages') else doc:
                    full_text += page.extract_text() if hasattr(page, 'extract_text') else page.get_text()
                table_sections = []

            # RFP 섹션 추출
            rfp_sections = await self.proposal_generator.rfp_extract_sections(full_text)
            
            # 제안서 상태 초기화
            proposal_data = {
                "id": proposal_id,
                "filename": filename,
                "created_at": datetime.now().isoformat(),
                "status": ProposalStatus.FILE_UPLOADED,
                "page_limit": page_limit,
                "rfp_sections": rfp_sections,
                "table_sections": table_sections,
                "full_text": full_text,
                "current_step": "analysis"
            }
            
            # 저장
            self.active_proposals[proposal_id] = proposal_data
            self._save_proposal(proposal_id, proposal_data)
            
            return proposal_id
            
        except Exception as e:
            logger.error(f"제안서 생성 중 오류 발생: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    def _save_proposal(self, proposal_id: str, data: dict):
        file_path = os.path.join(self.storage_dir, f"{proposal_id}.json")
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def get_proposal(self, proposal_id: str) -> Optional[dict]:
        if proposal_id in self.active_proposals:
            return self.active_proposals[proposal_id]
            
        file_path = os.path.join(self.storage_dir, f"{proposal_id}.json")
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                self.active_proposals[proposal_id] = data
                return data
        except FileNotFoundError:
            return None

    async def update_proposal_status(self, proposal_id: str, status: ProposalStatus):
        proposal = self.get_proposal(proposal_id)
        if proposal:
            proposal['status'] = status
            self._save_proposal(proposal_id, proposal)
            return proposal
        return None

# 전역 객체 초기화
proposal_manager = ProposalManager()

@app.post("/api/upload")
async def upload_file(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    try:
        file_content = await file.read()
        proposal_id = await proposal_manager.create_proposal(file_content, file.filename)
        
        return JSONResponse({
            "status": "success",
            "proposal_id": proposal_id,
            "message": "파일이 성공적으로 업로드되었습니다."
        })
        
    except Exception as e:
        logger.error(f"파일 업로드 중 오류 발생: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/proposal/{proposal_id}")
async def get_proposal_status(proposal_id: str):
    proposal_data = proposal_manager.get_proposal(proposal_id)
    if not proposal_data:
        raise HTTPException(status_code=404, detail="제안서를 찾을 수 없습니다.")
    return JSONResponse(proposal_data)

@app.post("/api/proposal/{proposal_id}/analyze")
async def analyze_proposal(proposal_id: str):
    proposal_data = proposal_manager.get_proposal(proposal_id)
    if not proposal_data:
        raise HTTPException(status_code=404, detail="제안서를 찾을 수 없습니다.")
        
    try:
        # 분석 단계 진행
        await proposal_manager.update_proposal_status(
            proposal_id, 
            ProposalStatus.PLANNING_IN_PROGRESS
        )
        
        # 여기에 분석 로직 추가
        
        return JSONResponse({
            "status": "success",
            "message": "분석이 완료되었습니다.",
            "proposal_data": proposal_data
        })
        
    except Exception as e:
        logger.error(f"제안서 분석 중 오류 발생: {str(e)}")
        await proposal_manager.update_proposal_status(
            proposal_id, 
            ProposalStatus.ERROR
        )
        raise HTTPException(status_code=500, detail=str(e))
    
    
@app.post("/api/proposal/{proposal_id}/planning")
async def plan_proposal(proposal_id: str):
    """목차 생성 및 기획 단계"""
    proposal_data = proposal_manager.get_proposal(proposal_id)
    if not proposal_data:
        raise HTTPException(status_code=404, detail="제안서를 찾을 수 없습니다.")
        
    try:
        # 페이지 수 계산
        presentation_time = proposal_data.get("rfp_sections", {}).get("발표시간", "15분")
        page_limit = proposal_data.get("page_limit", 30)
        
        # 발표시간에서 숫자만 추출
        time_match = re.search(r'(\d+)', presentation_time)
        presentation_minutes = int(time_match.group(1)) if time_match else 15
        
        # 페이지 수 계산 (발표시간 * 2, 단 page_limit를 넘지 않도록)
        recommended_pages = min(presentation_minutes * 2, page_limit) if page_limit else presentation_minutes * 2
        
        # 목차 생성
        toc_data = await proposal_manager.proposal_generator.generate_toc(
            proposal_data["rfp_sections"],
            recommended_pages,
            presentation_minutes
        )
        
        # 제안서 데이터 업데이트
        proposal_data["toc_data"] = toc_data
        proposal_data["recommended_pages"] = recommended_pages
        proposal_data["status"] = ProposalStatus.PLANNING_COMPLETED
        
        proposal_manager._save_proposal(proposal_id, proposal_data)
        
        return JSONResponse({
            "status": "success",
            "message": "목차가 생성되었습니다.",
            "data": toc_data
        })
        
    except Exception as e:
        logger.error(f"목차 생성 중 오류 발생: {str(e)}")
        await proposal_manager.update_proposal_status(proposal_id, ProposalStatus.ERROR)
        raise HTTPException(status_code=500, detail=str(e))
    
@app.post("/api/proposal/{proposal_id}/generate")
async def generate_proposal_content(proposal_id: str):
    """목차를 기반으로 각 페이지의 내용을 생성"""
    proposal_data = proposal_manager.get_proposal(proposal_id)
    if not proposal_data:
        raise HTTPException(status_code=404, detail="제안서를 찾을 수 없습니다.")
    
    try:
        toc_data = proposal_data["toc_data"]
        rfp_data = proposal_data["rfp_sections"]
        
        # 각 섹션별로 내용 생성
        generated_content = []
        for section in toc_data["sections"]:
            section_content = await proposal_manager.proposal_generator.generate_page_content(
                section["title"],
                rfp_data,
                {"pages": section["pages"]}
            )
            generated_content.append({
                "section": section["title"],
                "content": section_content
            })
        
        # 결과 저장
        proposal_data["generated_content"] = generated_content
        proposal_data["status"] = ProposalStatus.COMPLETED
        proposal_manager._save_proposal(proposal_id, proposal_data)
        
        return JSONResponse({
            "status": "success",
            "message": "제안서 내용이 생성되었습니다.",
            "data": generated_content
        })
        
    except Exception as e:
        logger.error(f"제안서 내용 생성 중 오류 발생: {str(e)}")
        await proposal_manager.update_proposal_status(proposal_id, ProposalStatus.ERROR)
        raise HTTPException(status_code=500, detail=str(e))
    


# TOC 업데이트 엔드포인트
@app.post("/api/proposal/{proposal_id}/toc/update")
async def update_toc(
    proposal_id: str,
    update_data: dict
):
    try:
        proposal_data = proposal_manager.get_proposal(proposal_id)
        if not proposal_data:
            raise HTTPException(status_code=404, detail="제안서를 찾을 수 없습니다.")
        
        section_index = update_data.get('sectionIndex')
        subsection_index = update_data.get('subsectionIndex')
        new_title = update_data.get('newTitle')
        
        # toc_data가 없는 경우 체크
        if 'toc_data' not in proposal_data:
            raise HTTPException(status_code=400, detail="목차 데이터가 없습니다.")
            
        # 목차 데이터 업데이트
        toc_data = proposal_data['toc_data']
        if subsection_index is None:
            # 대목차 수정
            toc_data['sections'][int(section_index)]['title'] = new_title
        else:
            # 중목차 수정
            toc_data['sections'][int(section_index)]['subsections'][int(subsection_index)]['title'] = new_title
            
        # 변경사항 저장
        proposal_manager._save_proposal(proposal_id, proposal_data)
        
        return JSONResponse({
            "status": "success",
            "message": "목차가 업데이트되었습니다."
        })
        
    except Exception as e:
        logger.error(f"목차 업데이트 중 오류 발생: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# PPTX 다운로드 엔드포인트
@app.get("/api/proposal/{proposal_id}/download")
async def download_proposal(proposal_id: str):
    try:
        proposal_data = proposal_manager.get_proposal(proposal_id)
        if not proposal_data:
            raise HTTPException(status_code=404, detail="제안서를 찾을 수 없습니다.")
            
        # PPTX 생성
        prs = create_pptx(proposal_data)
        
        # 임시 파일로 저장
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            prs.save(tmp.name)
            
            return FileResponse(
                tmp.name,
                media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                filename=f'제안서_{proposal_id}.pptx'
            )
            
    except Exception as e:
        logger.error(f"PPTX 생성 중 오류 발생: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

def create_pptx(proposal_data: dict) -> Presentation:
    """제안서 데이터를 기반으로 PPTX 생성"""
    prs = Presentation()
    
    # 표지 슬라이드 생성
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "제안서"
    subtitle.text = proposal_data.get('rfp_sections', {}).get('사업개요', '').split('\n')[0]
    
    # 목차 슬라이드 생성
    toc_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(toc_slide_layout)
    title = slide.shapes.title
    title.text = "목 차"
    
    # 내용 슬라이드 생성
    content_slide_layout = prs.slide_layouts[2]
    
    for section in proposal_data.get('generated_content', []):
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        # 헤드카피를 제목으로
        title.text = section['content']['head_copy']
        
        # 서브카피와 내용을 본문에
        content_text = f"{section['content']['sub_copy']}\n\n{section['content']['content']}"
        content.text = content_text
        
        # 텍스트 크기 조정
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
    
    return prs    

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("prop:app", host="0.0.0.0", port=8005, reload=True)